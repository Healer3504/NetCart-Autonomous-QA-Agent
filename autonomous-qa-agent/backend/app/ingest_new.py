# backend/app/ingest_new.py

from pathlib import Path
from typing import Union, List, Dict
from fastapi import UploadFile
import pandas as pd
import os

try:
    import fitz  # PyMuPDF
except:
    fitz = None

try:
    from docx import Document
except:
    Document = None

try:
    from pptx import Presentation
except:
    Presentation = None

UPLOAD_DIR = Path("uploaded_files")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)


async def save_uploaded_file(file: UploadFile, dest_dir: Union[str, Path] = UPLOAD_DIR) -> Path:
    dest_dir = Path(dest_dir)
    dest_dir.mkdir(parents=True, exist_ok=True)

    file_path = dest_dir / file.filename
    content = await file.read()

    with open(file_path, "wb") as f:
        f.write(content)

    return file_path

def extract_text_from_path(path: Union[str, Path]) -> str:
    """
    Extract text from various file types including .md and .json
    """
    p = Path(path)
    ext = p.suffix.lower()

    try:
        if ext == ".pdf" and fitz:
            doc = fitz.open(str(p))
            return "\n".join([page.get_text() for page in doc])

        if ext == ".docx" and Document:
            return "\n".join([para.text for para in Document(str(p)).paragraphs])

        if ext in [".ppt", ".pptx"] and Presentation:
            ppt = Presentation(str(p))
            out = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        out.append(shape.text)
            return "\n".join(out)

        if ext == ".txt":
            return p.read_text(encoding="utf-8", errors="ignore")
        
        if ext == ".md":
            # Markdown files - read as plain text
            return p.read_text(encoding="utf-8", errors="ignore")
        
        if ext == ".json":
            # JSON files - read as plain text
            return p.read_text(encoding="utf-8", errors="ignore")

        if ext == ".csv":
            df = pd.read_csv(p)
            return "\n".join(df.astype(str).values.flatten())

        return "[UNSUPPORTED FILE TYPE]"
    except Exception as e:
        return f"[ERROR Extracting Text: {e}]"


def build_vector_store_from_texts(items: List[Dict]) -> int:
    try:
        import chromadb
        from sentence_transformers import SentenceTransformer
    except:
        return len(items)

    embeddings = SentenceTransformer("all-MiniLM-L6-v2")

    client = chromadb.PersistentClient(path="vector_store")
    collection = client.get_or_create_collection("docs")

    docs = [item["text"] for item in items]
    ids = [item["filename"] for item in items]
    vecs = embeddings.encode(docs).tolist()

    collection.upsert(ids=ids, embeddings=vecs, documents=docs)

    return len(items)
